"""
Index manager extracted from legacy single-file implementation.
"""

import os
import time
import threading
import concurrent.futures
from collections import deque
from pathlib import Path
import logging

from PySide6.QtCore import QObject, Signal

from ..constants import (
    LOG_DIR,
    IS_WINDOWS,
    SKIP_DIRS_LOWER,
    SKIP_EXTS,
)
from ..utils import should_skip_path, should_skip_dir, get_c_scan_dirs
from ..utils import format_size, format_time  # scoring utilities removed
from .dependencies import HAS_APSW, get_db_module
from .mft_scanner import enum_volume_files_mft
from .trigram_index import TrigramIndex

logger = logging.getLogger(__name__)

db_module = get_db_module()


class IndexManager(QObject):
    """索引管理器 - 管理文件索引数据库"""

    progress_signal = Signal(int, str)
    build_finished_signal = Signal()
    content_progress_signal = Signal(int, int, int, str)
    content_build_finished_signal = Signal(bool)
    fts_finished_signal = Signal()

    def __init__(self, db_path=None, config_mgr=None):
        super().__init__()
        self.config_mgr = config_mgr
        if db_path is None:
            idx_dir = LOG_DIR
            idx_dir.mkdir(exist_ok=True)
            self.db_path = str(idx_dir / "index.db")
        else:
            self.db_path = db_path

        self.conn = None
        self.lock = threading.RLock()
        self.is_ready = False
        self.is_building = False
        # flag to support cancelling content index build
        self._stop_content_build = False
        self.file_count = 0
        self.last_build_time = None
        self.last_build_duration = None
        self.has_fts = False
        self.used_mft = False

        # in-memory trigram index (prototype)
        try:
            self.trigram_index = TrigramIndex()
        except Exception:
            self.trigram_index = None

        self._init_db()

    def _init_db(self):
        """初始化数据库"""
        try:
            if HAS_APSW:
                self.conn = db_module.Connection(self.db_path)
            else:
                self.conn = db_module.connect(self.db_path, check_same_thread=False)

            cursor = self.conn.cursor()
            cursor.execute("PRAGMA journal_mode=WAL")
            cursor.execute("PRAGMA synchronous=NORMAL")
            cursor.execute("PRAGMA cache_size=-2000000")
            cursor.execute("PRAGMA temp_store=MEMORY")

            cursor.execute(
                """
                CREATE TABLE IF NOT EXISTS files (
                    id INTEGER PRIMARY KEY,
                    filename TEXT NOT NULL,
                    filename_lower TEXT NOT NULL,
                    full_path TEXT UNIQUE NOT NULL,
                    parent_dir TEXT NOT NULL,
                    extension TEXT,
                    size INTEGER DEFAULT 0,
                    mtime REAL DEFAULT 0,
                    is_dir INTEGER DEFAULT 0
                )
            """
            )

            try:
                fts_exists = False
                for row in cursor.execute(
                    "SELECT name FROM sqlite_master WHERE type='table' AND name='files_fts'"
                ):
                    fts_exists = True
                    break
                if not fts_exists:
                    cursor.execute(
                        "CREATE VIRTUAL TABLE files_fts USING fts5(filename, content=files, content_rowid=id)"
                    )
                    cursor.execute(
                        """
                        CREATE TRIGGER IF NOT EXISTS files_ai AFTER INSERT ON files BEGIN
                            INSERT INTO files_fts(rowid, filename) VALUES (new.id, new.filename);
                        END
                    """
                    )
                    cursor.execute(
                        """
                        CREATE TRIGGER IF NOT EXISTS files_ad AFTER DELETE ON files BEGIN
                            INSERT INTO files_fts(files_fts, rowid, filename) VALUES('delete', old.id, old.filename);
                        END
                    """
                    )
                # 尝试填充 FTS5 内容表（如果尚未填充）并检查行数一致性
                try:
                    cursor.execute("INSERT INTO files_fts(rowid, filename) SELECT id, filename FROM files")
                    if not HAS_APSW:
                        self.conn.commit()
                except Exception:
                    # 某些 SQLite/FTS 环境下 insert 可能失败；继续但记录日志
                    logger.debug("尝试填充 files_fts 失败（可能已存在或不支持 INSERT INTO fts），将跳过填充")
                try:
                    files_cnt = list(cursor.execute("SELECT COUNT(*) FROM files"))[0][0]
                    fts_cnt = list(cursor.execute("SELECT COUNT(*) FROM files_fts"))[0][0]
                    if files_cnt != fts_cnt:
                        logger.warning(f"FTS5 行数与主表不一致: files={files_cnt}, files_fts={fts_cnt}")
                except Exception:
                    logger.debug("无法比较 files 与 files_fts 的行数")
                self.has_fts = True
                logger.info("✅ FTS5 已启用")
                # 尝试创建内容 FTS 表（用于文件内容全文搜索）
                try:
                    content_exists = False
                    for row in cursor.execute(
                        "SELECT name FROM sqlite_master WHERE type='table' AND name='content_fts'"
                    ):
                        content_exists = True
                        break
                    if not content_exists:
                        # content 列为全文索引，path 与 fileid 不参与倒排索引
                        cursor.execute(
                            "CREATE VIRTUAL TABLE content_fts USING fts5(content, path UNINDEXED, fileid UNINDEXED)"
                        )
                    self.has_content_fts = True
                    logger.info("✅ content FTS5 表已创建（用于文件内容搜索）")
                except Exception as e:
                    self.has_content_fts = False
                    logger.warning(f"⚠️ content FTS5 不可用: {e}")
            except Exception as e:
                self.has_fts = False
                logger.warning(f"⚠️ FTS5 不可用: {e}")

            cursor.execute("CREATE INDEX IF NOT EXISTS idx_fn ON files(filename_lower)")
            cursor.execute("CREATE INDEX IF NOT EXISTS idx_parent ON files(parent_dir)")
            cursor.execute("CREATE INDEX IF NOT EXISTS idx_ext ON files(extension)")
            cursor.execute(
                "CREATE TABLE IF NOT EXISTS meta (key TEXT PRIMARY KEY, value TEXT)"
            )

            if not HAS_APSW:
                self.conn.commit()

            self._load_stats()
        except Exception as e:
            logger.error(f"❌ 数据库初始化错误: {e}")
            self.conn = None
            self.is_ready = False

    def _load_stats(self, preserve_mft=False):
        """加载统计信息"""
        if not self.conn:
            return
        try:
            with self.lock:
                cursor = self.conn.cursor()

                count_result = list(cursor.execute("SELECT COUNT(*) FROM files"))
                self.file_count = count_result[0][0] if count_result else 0

                time_row = list(cursor.execute("SELECT value FROM meta WHERE key='build_time'"))
                if time_row and time_row[0][0]:
                    try:
                        self.last_build_time = float(time_row[0][0])
                    except (ValueError, TypeError):
                        self.last_build_time = None
                else:
                    self.last_build_time = None

                dur_row = list(cursor.execute("SELECT value FROM meta WHERE key='build_duration'"))
                if dur_row and dur_row[0][0]:
                    try:
                        self.last_build_duration = float(dur_row[0][0])
                    except (ValueError, TypeError):
                        self.last_build_duration = None
                else:
                    self.last_build_duration = None

                if self.last_build_time in (None, 0):
                    # 回退：若未写入 meta，则使用索引库文件的修改时间
                    try:
                        if os.path.exists(self.db_path):
                            self.last_build_time = os.path.getmtime(self.db_path)
                    except Exception:
                        pass

                if not preserve_mft:
                    mft_row = list(cursor.execute("SELECT value FROM meta WHERE key='used_mft'"))
                    self.used_mft = bool(mft_row and mft_row[0][0] == "1")

            self.is_ready = self.file_count > 0
        except Exception as e:
            logger.error(f"加载统计信息失败: {e}")
            self.file_count = 0
            self.is_ready = False

    def reload_stats(self):
        if not self.is_building:
            self._load_stats(preserve_mft=True)

    def force_reload_stats(self):
        self._load_stats(preserve_mft=True)

    def close(self):
        with self.lock:
            if self.conn:
                try:
                    self.conn.close()
                    logger.info("数据库连接已关闭")
                except Exception as e:
                    logger.warning(f"关闭数据库连接失败: {e}")
                finally:
                    self.conn = None

    def search(self, keywords, scope_targets, limit=50000):
        if not self.conn or not self.is_ready:
            logger.warning("搜索请求时索引不可用或未准备好: conn=%s, is_ready=%s", bool(self.conn), self.is_ready)
            # 静默回退：不抛错，返回空结果，避免打断搜索流程
            return []

        try:
            with self.lock:
                cursor = self.conn.cursor()
                keyword_str = ' '.join(keywords) if isinstance(keywords, list) else keywords
                parsed_keywords, filters, or_keywords, not_keywords = self._parse_search_syntax(keyword_str)
                
                # 详细调试日志
                logger.info(f"🔍 搜索参数 - 原始关键词: {keywords}, 关键词字符串: {keyword_str}")
                logger.info(f"🔍 解析结果 - AND关键词: {parsed_keywords}, OR关键词: {or_keywords}, NOT关键词: {not_keywords}, 过滤器: {filters}")

                if not parsed_keywords and not or_keywords and not any([
                    filters['ext'], filters['ext_list'], filters['size_min'], filters['size_max'],
                    filters['dm_after'], filters['dm_before'], filters['type'], filters['path'],
                    filters['len_min'], filters['len_max']
                ]):
                    return []

                # If the configuration requests a simple 'Everything'-like
                # substring search mode, use a straightforward SQL query that
                # matches keywords against both filename and full_path. This
                # keeps behavior simple and predictable for users who prefer
                # Everything-style searching. Otherwise, fall back to using
                # the in-memory trigram index when available for candidate
                # selection, and finally the SQL LIKE path.
                simple_mode = True
                try:
                    if getattr(self, 'config_mgr', None) is not None:
                        simple_mode = bool(self.config_mgr.get_search_simple_mode())
                except Exception:
                    simple_mode = True

                if simple_mode:
                    # Everything 风格简单模式：匹配文件名或路径
                    match_on_path = True
                    conditions = []
                    params = []

                    # 通配符转换辅助函数
                    def wildcard_to_sql(pattern):
                        """将 Everything 风格通配符转换为 SQL LIKE 模式
                        * -> %
                        ? -> _
                        """
                        # 转义 SQL LIKE 特殊字符（除了即将替换的 * 和 ?）
                        pattern = pattern.replace('[', r'\[').replace('%', r'\%').replace('_', r'\_')
                        # 转换通配符
                        pattern = pattern.replace('*', '%').replace('?', '_')
                        return pattern

                    # AND 关键词（空格分隔）
                    for kw in parsed_keywords:
                        sql_pattern = wildcard_to_sql(kw)
                        conditions.append("(filename_lower LIKE ? ESCAPE '\\' OR lower(full_path) LIKE ? ESCAPE '\\')")
                        params.append(f"%{sql_pattern}%")
                        params.append(f"%{sql_pattern}%")

                    # OR 关键词（| 分隔）
                    if or_keywords:
                        or_conditions = []
                        for or_kw in or_keywords:
                            sql_pattern = wildcard_to_sql(or_kw)
                            or_conditions.append("(filename_lower LIKE ? ESCAPE '\\' OR lower(full_path) LIKE ? ESCAPE '\\')")
                            params.append(f"%{sql_pattern}%")
                            params.append(f"%{sql_pattern}%")
                        if or_conditions:
                            conditions.append(f"({' OR '.join(or_conditions)})")

                    # NOT 关键词（! 前缀）
                    for not_kw in not_keywords:
                        sql_pattern = wildcard_to_sql(not_kw)
                        conditions.append("NOT (filename_lower LIKE ? ESCAPE '\\' OR lower(full_path) LIKE ? ESCAPE '\\')")
                        params.append(f"%{sql_pattern}%")
                        params.append(f"%{sql_pattern}%")

                    # 扩展名过滤（支持多个扩展名 OR）
                    if filters['ext_list']:
                        ext_conditions = []
                        for ext in filters['ext_list']:
                            ext_conditions.append("extension = ?")
                            params.append(ext)
                        conditions.append(f"({' OR '.join(ext_conditions)})")
                    elif filters['ext']:
                        conditions.append("extension = ?")
                        params.append(filters['ext'])

                    # 文件类型过滤
                    if filters['type'] == 'folder':
                        conditions.append("is_dir = 1")
                    elif filters['type'] == 'file':
                        conditions.append("is_dir = 0")

                    # 大小过滤
                    if filters['size_min'] > 0:
                        conditions.append("size > ?")
                        params.append(filters['size_min'])
                    if filters['size_max'] > 0:
                        conditions.append("size < ?")
                        params.append(filters['size_max'])

                    # 修改时间过滤
                    if filters['dm_after'] > 0:
                        conditions.append("(mtime >= ? OR mtime = 0)")
                        params.append(filters['dm_after'])
                    if filters['dm_before'] > 0:
                        conditions.append("(mtime <= ? OR mtime = 0)")
                        params.append(filters['dm_before'])

                    # 路径长度过滤
                    if filters['len_min'] > 0:
                        conditions.append("LENGTH(full_path) > ?")
                        params.append(filters['len_min'])
                    if filters['len_max'] > 0:
                        conditions.append("LENGTH(full_path) < ?")
                        params.append(filters['len_max'])

                    # 路径包含过滤
                    if filters['path']:
                        conditions.append("lower(full_path) LIKE ?")
                        params.append(f"%{filters['path']}%")

                    # 内容搜索（使用 content_fts FTS5 表）
                    if filters.get('content'):
                        if getattr(self, 'has_content_fts', False):
                            conditions.append("EXISTS (SELECT 1 FROM content_fts c WHERE c.rowid = files.id AND c.content MATCH ?)")
                            params.append(filters['content'])
                        else:
                            # 如果 FTS5 内容索引不可用，记录并忽略（或可选择回退到慢速扫描）
                            logger.warning("内容搜索被请求但 content_fts 不可用，跳过内容过滤")

                    where_clause = " AND ".join(conditions) if conditions else "1=1"
                    sql = f"""
                        SELECT filename, full_path, size, mtime, is_dir
                        FROM files
                        WHERE {where_clause}
                        LIMIT ?
                    """
                    params.append(limit)
                    
                    logger.info(f"🔍 SQL查询 - simple_mode: {simple_mode}, match_on_path: {match_on_path}")
                    logger.info(f"🔍 SQL: {sql}")
                    logger.info(f"🔍 参数: {params}")
                    
                    raw_results = list(cursor.execute(sql, tuple(params)))
                    logger.info(f"🔍 SQL返回原始结果数: {len(raw_results)}")
                else:
                    # 高级模式：只匹配文件名
                    match_on_path = False
                    conditions = []
                    params = []

                    # 通配符转换辅助函数
                    def wildcard_to_sql(pattern):
                        """将 Everything 风格通配符转换为 SQL LIKE 模式"""
                        pattern = pattern.replace('[', r'\[').replace('%', r'\%').replace('_', r'\_')
                        pattern = pattern.replace('*', '%').replace('?', '_')
                        return pattern

                    # AND 关键词
                    for kw in parsed_keywords:
                        sql_pattern = wildcard_to_sql(kw)
                        conditions.append("filename_lower LIKE ? ESCAPE '\\'")
                        params.append(f"%{sql_pattern}%")

                    # OR 关键词
                    if or_keywords:
                        or_conditions = []
                        for or_kw in or_keywords:
                            sql_pattern = wildcard_to_sql(or_kw)
                            or_conditions.append("filename_lower LIKE ? ESCAPE '\\'")
                            params.append(f"%{sql_pattern}%")
                        if or_conditions:
                            conditions.append(f"({' OR '.join(or_conditions)})")

                    # NOT 关键词
                    for not_kw in not_keywords:
                        sql_pattern = wildcard_to_sql(not_kw)
                        conditions.append("NOT filename_lower LIKE ? ESCAPE '\\'")
                        params.append(f"%{sql_pattern}%")

                    # 扩展名过滤（支持多个扩展名 OR）
                    if filters['ext_list']:
                        ext_conditions = []
                        for ext in filters['ext_list']:
                            ext_conditions.append("extension = ?")
                            params.append(ext)
                        conditions.append(f"({' OR '.join(ext_conditions)})")
                    elif filters['ext']:
                        conditions.append("extension = ?")
                        params.append(filters['ext'])

                    # 文件类型过滤
                    if filters['type'] == 'folder':
                        conditions.append("is_dir = 1")
                    elif filters['type'] == 'file':
                        conditions.append("is_dir = 0")

                    # 大小过滤
                    if filters['size_min'] > 0:
                        conditions.append("size > ?")
                        params.append(filters['size_min'])
                    if filters['size_max'] > 0:
                        conditions.append("size < ?")
                        params.append(filters['size_max'])

                    # 修改时间过滤
                    if filters['dm_after'] > 0:
                        conditions.append("(mtime >= ? OR mtime = 0)")
                        params.append(filters['dm_after'])
                    if filters['dm_before'] > 0:
                        conditions.append("(mtime <= ? OR mtime = 0)")
                        params.append(filters['dm_before'])

                    # 路径长度过滤
                    if filters['len_min'] > 0:
                        conditions.append("LENGTH(full_path) > ?")
                        params.append(filters['len_min'])
                    if filters['len_max'] > 0:
                        conditions.append("LENGTH(full_path) < ?")
                        params.append(filters['len_max'])

                    # 路径包含过滤
                    if filters['path']:
                        conditions.append("lower(full_path) LIKE ?")
                        params.append(f"%{filters['path']}%")

                    # 内容搜索（使用 content_fts FTS5 表）
                    if filters.get('content'):
                        if getattr(self, 'has_content_fts', False):
                            conditions.append("EXISTS (SELECT 1 FROM content_fts c WHERE c.rowid = files.id AND c.content MATCH ?)")
                            params.append(filters['content'])
                        else:
                            logger.warning("内容搜索被请求但 content_fts 不可用，跳过内容过滤")

                    where_clause = " AND ".join(conditions) if conditions else "1=1"
                    sql = f"""
                        SELECT filename, full_path, size, mtime, is_dir
                        FROM files
                        WHERE {where_clause}
                        LIMIT ?
                    """
                    params.append(limit)
                    
                    logger.info(f"🔍 SQL查询 - simple_mode: {simple_mode}, match_on_path: {match_on_path}")
                    logger.info(f"🔍 SQL: {sql}")
                    logger.info(f"🔍 参数: {params}")
                    
                    raw_results = list(cursor.execute(sql, tuple(params)))
                    logger.info(f"🔍 SQL返回原始结果数: {len(raw_results)}")
                
                # Apply scope filtering and path/dir skip logic
                scope_drives = set()
                scope_paths = []
                if scope_targets:
                    for t in scope_targets:
                        t_norm = os.path.normpath(t).lower().rstrip("\\")
                        drv = Path(t_norm).drive.lower()
                        if drv and (t_norm == drv or t_norm == drv + "\\"):
                            scope_drives.add(drv)
                        else:
                            scope_paths.append(t_norm)

                filtered = []
                for fn, fp, sz, mt, is_dir in raw_results:
                    path_norm = os.path.normpath(fp)
                    path_lower = path_norm.lower()
                    path_drive = Path(path_lower).drive.lower()

                    if scope_targets:
                        ok = False
                        if scope_drives and path_drive in scope_drives:
                            ok = True
                        if not ok and scope_paths:
                            for p in scope_paths:
                                if path_lower == p or path_lower.startswith(p + "\\"):
                                    ok = True
                                    break
                        if not ok:
                            continue

                    if filters['path'] and filters['path'] not in path_lower:
                        continue

                    if should_skip_path(path_lower):
                        continue

                    name_lower = fn.lower()
                    if is_dir:
                        if should_skip_dir(name_lower, path_lower):
                            continue
                    else:
                        if os.path.splitext(name_lower)[1] in SKIP_EXTS:
                            continue

                    filtered.append((fn, fp, sz, mt, is_dir))

            logger.info(f"🔍 范围/路径过滤后结果数: {len(filtered)}")
            
            if filtered and filters.get('dm_after', 0) > 0:
                needs_fix_count = sum(1 for item in filtered if item[3] == 0)
                if needs_fix_count > 0:
                    logger.info(f"⚠️ 发现 {needs_fix_count} 个文件 mtime=0，使用快速补齐（多线程）...")
                    start_time = time.time()
                    max_fix = 10000
                    needs_fix = [(i, item[1]) for i, item in enumerate(filtered) if item[3] == 0][:max_fix]

                    def get_mtime_fast(idx_path):
                        idx, fpath = idx_path
                        try:
                            return (idx, os.stat(fpath).st_mtime, fpath)
                        except Exception:
                            return (idx, 0, fpath)

                    fixed_items = {}
                    db_updates = []

                    with concurrent.futures.ThreadPoolExecutor(max_workers=16) as executor:
                        for idx, mt, fpath in executor.map(get_mtime_fast, needs_fix):
                            fixed_items[idx] = mt
                            if mt > 0:
                                db_updates.append((mt, fpath))

                    new_filtered = []
                    for i, (fn, fp, sz, mt, is_dir) in enumerate(filtered):
                        if i in fixed_items:
                            mt = fixed_items[i]
                        if mt > 0 and mt >= filters['dm_after']:
                            new_filtered.append((fn, fp, sz, mt, is_dir))

                    filtered = new_filtered
                    elapsed = time.time() - start_time
                    logger.info(f"✅ 补齐完成: {len(needs_fix)} 个文件，耗时 {elapsed:.2f}s，剩余 {len(filtered)} 个")

                    if db_updates:
                        def update_db():
                            try:
                                with self.lock:
                                    cursor = self.conn.cursor()
                                    cursor.executemany("UPDATE files SET mtime=? WHERE full_path=?", db_updates)
                                    if not HAS_APSW:
                                        self.conn.commit()
                                logger.info(f"📝 已缓存 {len(db_updates)} 个文件的 mtime 到数据库")
                            except Exception as e:
                                logger.debug(f"数据库更新失败: {e}")

                        threading.Thread(target=update_db, daemon=True).start()

            logger.info(f"✅ 搜索完成，最终返回结果数: {len(filtered)}")
            return filtered

        except Exception as e:
            logger.error(f"搜索错误: {e}")
            import traceback
            traceback.print_exc()
            # 静默回退：异常时返回空列表，避免前端弹窗
            return []

    def _search_like(self, cursor, keywords, limit):
        wheres = " AND ".join(["filename_lower LIKE ?"] * len(keywords))
        sql = f"""
            SELECT filename, full_path, size, mtime, is_dir
            FROM files
            WHERE {wheres}
            LIMIT ?
        """
        params = tuple([f"%{kw}%" for kw in keywords] + [limit])
        return list(cursor.execute(sql, params))

    def get_stats(self):
        self._load_stats(preserve_mft=True)
        return {
            "count": self.file_count,
            "ready": self.is_ready,
            "building": self.is_building,
            "time": self.last_build_time,
            "duration": self.last_build_duration,
            "path": self.db_path,
            "has_fts": self.has_fts,
            "used_mft": self.used_mft,
        }

    def build_index(self, drives, stop_fn=None):
        from . import mft_scanner  # avoid circular

        if not self.conn or self.is_building:
            return

        self.is_building = True
        self.is_ready = False
        self.used_mft = False
        mft_scanner.MFT_AVAILABLE = False
        build_start = time.time()

        try:
            logger.info("🚀 开始构建索引...")
            self.progress_signal.emit(0, "阶段1/5: 清理旧数据...")

            with self.lock:
                cursor = self.conn.cursor()
                cursor.execute("DROP TRIGGER IF EXISTS files_ai")
                cursor.execute("DROP TRIGGER IF EXISTS files_ad")
                cursor.execute("DROP TABLE IF EXISTS files_fts")
                cursor.execute("DROP TABLE IF EXISTS files")
                cursor.execute(
                    """
                    CREATE TABLE files (
                        id INTEGER PRIMARY KEY,
                        filename TEXT NOT NULL,
                        filename_lower TEXT NOT NULL,
                        full_path TEXT UNIQUE NOT NULL,
                        parent_dir TEXT NOT NULL,
                        extension TEXT,
                        size INTEGER DEFAULT 0,
                        mtime REAL DEFAULT 0,
                        is_dir INTEGER DEFAULT 0
                    )
                """
                )
                if not HAS_APSW:
                    self.conn.commit()
                self.has_fts = False
                self.file_count = 0

            logger.info(f"✅ 阶段1完成: {time.time() - build_start:.2f}s")

            self.progress_signal.emit(0, "阶段2/5: MFT扫描...")
            all_drives = [d.upper().rstrip(":\\") for d in drives if os.path.exists(d)]
            c_allowed_paths = get_c_scan_dirs(self.config_mgr)
            all_data = []
            failed_drives = []

            if all_drives and IS_WINDOWS:
                data_lock = threading.Lock()

                def scan_one(drv):
                    try:
                        allowed = c_allowed_paths if drv == "C" else None
                        data = enum_volume_files_mft(
                            drv, SKIP_DIRS_LOWER, SKIP_EXTS, allowed_paths=allowed
                        )
                        with data_lock:
                            all_data.extend(data)
                        return drv, len(data)
                    except Exception as e:
                        logger.error(f"扫描驱动器 {drv} 失败: {e}")
                        return drv, -1

                with concurrent.futures.ThreadPoolExecutor(
                    max_workers=min(len(all_drives), 4)
                ) as ex:
                    futures = [ex.submit(scan_one, d) for d in all_drives]
                    for future in concurrent.futures.as_completed(futures):
                        if stop_fn and stop_fn():
                            break
                        drv, result = future.result()
                        if result < 0:
                            failed_drives.append(drv)
                        self.progress_signal.emit(
                            len(all_data),
                            f"MFT {drv}: {result if result >= 0 else '失败'}",
                        )

                if all_data:
                    self.used_mft = True

            logger.info(
                f"✅ 阶段2完成: {time.time() - build_start:.2f}s, 扫描到 {len(all_data):,} 条"
            )

            if all_data:
                self.progress_signal.emit(len(all_data), "阶段3/5: 写入数据库...")
                write_start = time.time()

                with self.lock:
                    cursor = self.conn.cursor()
                    cursor.execute("PRAGMA synchronous=OFF")
                    cursor.execute("PRAGMA journal_mode=MEMORY")
                    cursor.execute("PRAGMA locking_mode=EXCLUSIVE")
                    cursor.execute("PRAGMA temp_store=MEMORY")
                    cursor.execute("PRAGMA cache_size=-500000")
                    cursor.execute("PRAGMA mmap_size=268435456")

                    if HAS_APSW:
                        with self.conn:
                            cursor.executemany(
                                "INSERT OR IGNORE INTO files VALUES(NULL,?,?,?,?,?,?,?,?)",
                                all_data
                            )
                    else:
                        cursor.executemany(
                            "INSERT OR IGNORE INTO files VALUES(NULL,?,?,?,?,?,?,?,?)",
                            all_data
                        )
                        self.conn.commit()

                self.file_count = len(all_data)
                write_time = time.time() - write_start
                logger.info(f"✅ 阶段3完成: {write_time:.2f}s, 写入 {len(all_data):,} 条")

                # build in-memory trigram index from all_data for fast candidate selection
                try:
                    if getattr(self, 'trigram_index', None) is not None:
                        docs = []
                        for fn, fn_lower, fp, cur, ext, sz, mt, is_dir in all_data:
                            docs.append({
                                'filename': fn,
                                'dir_path': cur,
                                'fullpath': fp,
                                'size': sz,
                                'mtime': mt,
                                'type_code': is_dir,
                            })
                        # build index (in-memory)
                        self.trigram_index.build_index(docs)
                except Exception:
                    pass

            self.progress_signal.emit(self.file_count, "阶段4/5: 创建索引...")
            with self.lock:
                cursor = self.conn.cursor()
                cursor.execute("CREATE INDEX IF NOT EXISTS idx_fn ON files(filename_lower)")
                cursor.execute("CREATE INDEX IF NOT EXISTS idx_parent ON files(parent_dir)")
                cursor.execute("PRAGMA synchronous=NORMAL")
                cursor.execute("PRAGMA journal_mode=WAL")
                now_ts = time.time()
                cursor.execute(
                    "INSERT OR REPLACE INTO meta (key, value) VALUES ('build_time', ?)",
                    (str(now_ts),),
                )
                cursor.execute(
                    "INSERT OR REPLACE INTO meta (key, value) VALUES ('build_duration', ?)",
                    (str(time.time() - build_start),),
                )
                cursor.execute(
                    "INSERT OR REPLACE INTO meta (key, value) VALUES ('used_mft', ?)",
                    ("1" if self.used_mft else "0",),
                )
                if not HAS_APSW:
                    self.conn.commit()

            logger.info(f"✅ 阶段4完成: {time.time() - build_start:.2f}s")

            self.progress_signal.emit(self.file_count, "阶段5/5: 构建全文索引(后台)...")

            def build_fts_async():
                try:
                    logger.info("📝 后台构建 FTS5...")
                    fts_start = time.time()
                    with self.lock:
                        cursor = self.conn.cursor()
                        cursor.execute(
                            "CREATE VIRTUAL TABLE IF NOT EXISTS files_fts USING fts5(filename, content=files, content_rowid=id)"
                        )
                        cursor.execute("INSERT INTO files_fts(files_fts) VALUES('rebuild')")
                        cursor.execute(
                            """
                            CREATE TRIGGER IF NOT EXISTS files_ai AFTER INSERT ON files BEGIN
                                INSERT INTO files_fts(rowid, filename) VALUES (new.id, new.filename);
                            END
                        """
                        )
                        cursor.execute(
                            """
                            CREATE TRIGGER IF NOT EXISTS files_ad AFTER DELETE ON files BEGIN
                                INSERT INTO files_fts(files_fts, rowid, filename) VALUES('delete', old.id, old.filename);
                            END
                        """
                        )
                        if not HAS_APSW:
                            self.conn.commit()
                        self.has_fts = True
                    logger.info(f"✅ FTS5 构建完成: {time.time() - fts_start:.2f}s")
                except Exception as e:
                    logger.warning(f"⚠️ FTS5 构建失败: {e}")
                    self.has_fts = False
                self.fts_finished_signal.emit()

            threading.Thread(target=build_fts_async, daemon=True).start()

            for drv in failed_drives:
                if stop_fn and stop_fn():
                    break
                paths_to_scan = c_allowed_paths if drv == "C" else [f"{drv}:\\"]
                for path in paths_to_scan:
                    logger.info(f"[传统扫描] {path}")
                    self._scan_dir(path, c_allowed_paths if drv == "C" else None, stop_fn)

            try:
                with self.lock:
                    cursor = self.conn.cursor()
                    final_count = list(cursor.execute("SELECT COUNT(*) FROM files"))[0][0]
                    self.file_count = final_count
            except Exception:
                pass

            total_time = time.time() - build_start
            logger.info(f"✅ 索引构建完成: {self.file_count:,} 条, 总耗时 {total_time:.2f}s")
            self.is_ready = self.file_count > 0
            self.build_finished_signal.emit()

        except Exception as e:
            import traceback
            logger.error(f"❌ 构建错误: {e}")
            traceback.print_exc()
        finally:
            self.is_building = False

    def _scan_dir(self, target, allowed_paths=None, stop_fn=None):
        try:
            if not os.path.exists(target):
                return
        except (OSError, PermissionError):
            logger.warning(f"无法访问目录: {target}")
            return

        allowed_paths_lower = (
            [p.lower().rstrip("\\") for p in allowed_paths] if allowed_paths else None
        )
        batch = []
        stack = deque([target])

        while stack:
            if stop_fn and stop_fn():
                break
            cur = stack.pop()
            if should_skip_path(cur.lower(), allowed_paths_lower):
                continue
            try:
                with os.scandir(cur) as it:
                    for e in it:
                        if stop_fn and stop_fn():
                            break
                        if not e.name or e.name.startswith((".", "$")):
                            continue
                        try:
                            is_dir = e.is_dir()
                            st = e.stat(follow_symlinks=False)
                        except (OSError, PermissionError):
                            continue

                        path_lower = e.path.lower()
                        if is_dir:
                            if should_skip_dir(e.name.lower(), path_lower, allowed_paths_lower):
                                continue
                            stack.append(e.path)
                            batch.append((e.name, e.name.lower(), e.path, cur, "", 0, st.st_mtime, 1))
                        else:
                            ext = os.path.splitext(e.name)[1].lower()
                            if ext in SKIP_EXTS:
                                continue
                            batch.append((e.name, e.name.lower(), e.path, cur, ext, st.st_size, st.st_mtime, 0))

                        if len(batch) >= 20000:
                            with self.lock:
                                cursor = self.conn.cursor()
                                cursor.executemany(
                                    "INSERT OR IGNORE INTO files VALUES(NULL,?,?,?,?,?,?,?,?)",
                                    batch,
                                )
                                if not HAS_APSW:
                                    self.conn.commit()
                                # best-effort: add to in-memory trigram index
                                try:
                                    if getattr(self, 'trigram_index', None) is not None:
                                        for fn, fn_lower, fp, curp, ext, sz, mt, is_dir in batch:
                                            doc = {
                                                'filename': fn,
                                                'dir_path': curp,
                                                'fullpath': fp,
                                                'size': sz,
                                                'mtime': mt,
                                                'type_code': is_dir,
                                            }
                                            try:
                                                self.trigram_index.add_doc(doc)
                                            except Exception:
                                                pass
                                except Exception:
                                    pass
                            self.file_count += len(batch)
                            self.progress_signal.emit(self.file_count, cur)
                            batch = []
            except (PermissionError, OSError):
                continue

        if batch:
            with self.lock:
                cursor = self.conn.cursor()
                cursor.executemany(
                    "INSERT OR IGNORE INTO files VALUES(NULL,?,?,?,?,?,?,?,?)", batch
                )
                if not HAS_APSW:
                    self.conn.commit()
                # best-effort: add to in-memory trigram index
                try:
                    if getattr(self, 'trigram_index', None) is not None:
                        for fn, fn_lower, fp, curp, ext, sz, mt, is_dir in batch:
                            doc = {
                                'filename': fn,
                                'dir_path': curp,
                                'fullpath': fp,
                                'size': sz,
                                'mtime': mt,
                                'type_code': is_dir,
                            }
                            try:
                                self.trigram_index.add_doc(doc)
                            except Exception:
                                pass
                except Exception:
                    pass
            self.file_count += len(batch)

    def rebuild_drive(self, drive_letter, progress_callback=None, stop_fn=None):
        if not self.conn:
            return
        if self.is_building:
            logger.warning("索引正在构建中，跳过")
            return

        self.is_building = True
        drive = drive_letter.upper().rstrip(":\\")
        try:
            logger.info(f"开始重建 {drive}: 盘索引...")
            with self.lock:
                cursor = self.conn.cursor()
                cursor.execute("DELETE FROM files WHERE full_path LIKE ?", (f"{drive}:%",))
                if not HAS_APSW:
                    self.conn.commit()

            c_allowed_paths = get_c_scan_dirs(self.config_mgr)
            allowed_paths = c_allowed_paths if drive == 'C' else None

            try:
                data = enum_volume_files_mft(drive, SKIP_DIRS_LOWER, SKIP_EXTS, allowed_paths)
                if data:
                    logger.info(f"开始写入 {len(data)} 条记录...")
                    write_start = time.time()
                    with self.lock:
                        cursor = self.conn.cursor()
                        cursor.execute("PRAGMA synchronous=OFF")
                        cursor.execute("PRAGMA journal_mode=OFF")
                        cursor.execute("PRAGMA locking_mode=EXCLUSIVE")
                        cursor.execute("PRAGMA temp_store=MEMORY")
                        cursor.execute("PRAGMA cache_size=-500000")

                        if HAS_APSW:
                            with self.conn:
                                cursor.executemany(
                                    "INSERT OR IGNORE INTO files VALUES(NULL,?,?,?,?,?,?,?,?)",
                                    data
                                )
                                now_ts = time.time()
                                cursor.execute(
                                    "INSERT OR REPLACE INTO meta (key, value) VALUES ('build_time', ?)",
                                    (str(now_ts),)
                                )
                                cursor.execute(
                                    "INSERT OR REPLACE INTO meta (key, value) VALUES ('build_duration', ?)",
                                    (str(time.time() - build_start),)
                                )
                                cursor.execute(
                                    "INSERT OR REPLACE INTO meta (key, value) VALUES ('used_mft', '1')"
                                )
                        else:
                            cursor.execute("BEGIN TRANSACTION")
                            cursor.executemany(
                                "INSERT OR IGNORE INTO files VALUES(NULL,?,?,?,?,?,?,?,?)",
                                data
                            )
                            now_ts = time.time()
                            cursor.execute(
                                "INSERT OR REPLACE INTO meta (key, value) VALUES ('build_time', ?)",
                                (str(now_ts),)
                            )
                            cursor.execute(
                                "INSERT OR REPLACE INTO meta (key, value) VALUES ('build_duration', ?)",
                                (str(time.time() - build_start),)
                            )
                            cursor.execute(
                                "INSERT OR REPLACE INTO meta (key, value) VALUES ('used_mft', '1')"
                            )
                            cursor.execute("COMMIT")

                        cursor.execute("PRAGMA synchronous=NORMAL")
                        cursor.execute("PRAGMA journal_mode=WAL")
                        cursor.execute("PRAGMA locking_mode=NORMAL")

                    write_time = time.time() - write_start
                    logger.info(f"✅ {drive}: 盘索引重建完成，写入 {len(data)} 条记录，耗时 {write_time:.2f}s")

            except Exception as e:
                logger.error(f"MFT扫描失败: {e}")
                import traceback
                traceback.print_exc()

            self._load_stats(preserve_mft=True)
            self.is_ready = self.file_count > 0

        except Exception as e:
            logger.error(f"重建驱动器 {drive} 失败: {e}")
            import traceback
            traceback.print_exc()
        finally:
            self.is_building = False
            logger.info(f"{drive}: 盘索引重建流程结束")
            self.build_finished_signal.emit()

    def _parse_search_syntax(self, keyword_str):
        """解析 Everything 风格的搜索语法
        
        支持：
        - 通配符：* 和 ?
        - 布尔运算：| (OR), ! (NOT), 空格 (AND)
        - 过滤器：ext:, size:, dm:, folder:, file:, path:, len:, attrib:, datemodified:
        """
        import re
        import datetime

        keywords = []
        or_keywords = []  # OR 关键词列表
        not_keywords = []  # NOT 关键词列表
        filters = {
            'ext': None,
            'ext_list': [],  # 支持多个扩展名（OR）
            'size_min': 0,
            'size_max': 0,
            'dm_after': 0,
            'dm_before': 0,
            'type': None,
            'path': None,
            'content': None,
            'len_min': 0,
            'len_max': 0,
            'attrib_hidden': None,
            'attrib_readonly': None,
        }

        # 预处理：处理 | 分隔的 OR 语法
        # 例如：jpg|png|gif -> 会被拆分为多个 OR 选项
        def split_or_tokens(text):
            """分割 OR 表达式，支持 word1|word2|word3"""
            if '|' in text:
                return [t.strip() for t in text.split('|') if t.strip()]
            return [text]

        tokens = keyword_str.split()
        for token in tokens:
            token_lower = token.lower()
            
            # 处理 NOT 运算符（!关键词）
            if token.startswith('!') and len(token) > 1:
                not_term = token[1:].lower()
                # 支持通配符转换为正则
                not_keywords.append(not_term)
                continue
            
            # 扩展名过滤：ext:jpg 或 ext:jpg|png|gif
            if token_lower.startswith('ext:'):
                ext_part = token[4:].strip()
                for ext in split_or_tokens(ext_part):
                    if ext and not ext.startswith('.'):
                        ext = '.' + ext
                    filters['ext_list'].append(ext.lower())
                if filters['ext_list']:
                    filters['ext'] = filters['ext_list'][0]  # 兼容旧逻辑
                continue
            
            # 大小过滤：size:>1mb, size:<500kb, size:1mb..10mb
            if token_lower.startswith('size:'):
                size_part = token[5:].strip().lower()
                # 范围语法：size:1mb..10mb
                if '..' in size_part:
                    try:
                        min_str, max_str = size_part.split('..')
                        filters['size_min'] = self._parse_size(min_str)
                        filters['size_max'] = self._parse_size(max_str)
                    except:
                        pass
                else:
                    match = re.match(r'([<>])(\d+)(kb|mb|gb)?', size_part)
                    if match:
                        op, num, unit = match.groups()
                        num = int(num)
                        multiplier = {'kb': 1024, 'mb': 1024**2, 'gb': 1024**3}.get(unit, 1)
                        size_bytes = num * multiplier
                        if op == '>':
                            filters['size_min'] = size_bytes
                        else:
                            filters['size_max'] = size_bytes
                continue
            
            # 修改时间过滤：dm:today, dm:7d, dm:2024-12-01, dm:2024-12-01..2024-12-22
            if token_lower.startswith('dm:') or token_lower.startswith('datemodified:'):
                dm_part = token.split(':', 1)[1].strip().lower()
                now = time.time()
                day = 86400
                
                # 范围语法：dm:2024-12-01..2024-12-22
                if '..' in dm_part:
                    try:
                        start_str, end_str = dm_part.split('..')
                        filters['dm_after'] = self._parse_date(start_str)
                        filters['dm_before'] = self._parse_date(end_str) + day  # 包含结束日期
                    except:
                        pass
                elif dm_part == 'today':
                    today_start = datetime.datetime.now().replace(
                        hour=0, minute=0, second=0, microsecond=0
                    ).timestamp()
                    filters['dm_after'] = today_start
                elif dm_part.endswith('d') and dm_part[:-1].isdigit():
                    days = int(dm_part[:-1])
                    filters['dm_after'] = now - (days * day)
                elif dm_part.endswith('h') and dm_part[:-1].isdigit():
                    hours = int(dm_part[:-1])
                    filters['dm_after'] = now - (hours * 3600)
                elif re.match(r'\d{4}-\d{2}-\d{2}', dm_part):
                    # 精确日期：dm:2024-12-22
                    try:
                        filters['dm_after'] = self._parse_date(dm_part)
                    except:
                        pass
                continue
            
            # 路径长度过滤：len:>100, len:<50
            if token_lower.startswith('len:'):
                len_part = token[4:].strip()
                match = re.match(r'([<>])(\d+)', len_part)
                if match:
                    op, num = match.groups()
                    if op == '>':
                        filters['len_min'] = int(num)
                    else:
                        filters['len_max'] = int(num)
                continue
            
            # 文件属性过滤：attrib:h (hidden), attrib:r (readonly)
            if token_lower.startswith('attrib:'):
                attrib = token[7:].strip().lower()
                if 'h' in attrib:
                    filters['attrib_hidden'] = True
                if 'r' in attrib:
                    filters['attrib_readonly'] = True
                continue
            
            # 文件夹/文件类型过滤
            if token_lower.startswith('folder:'):
                filters['type'] = 'folder'
                rest = token[7:].strip()
                if rest:
                    keywords.append(rest.lower())
                continue
            
            if token_lower.startswith('file:'):
                filters['type'] = 'file'
                rest = token[5:].strip()
                if rest:
                    keywords.append(rest.lower())
                continue
            
            # 路径包含过滤
            if token_lower.startswith('path:'):
                path_part = token[5:].strip()
                if path_part:
                    filters['path'] = path_part.lower()
                continue
            # 文件内容搜索过滤器
            if token_lower.startswith('content:'):
                # 支持 content:"phrase with spaces" 或 content:word
                content_part = token.split(':', 1)[1].strip()
                if content_part:
                    filters['content'] = content_part
                continue
            
            # 处理 OR 关键词（包含 | 的）
            if '|' in token:
                or_keywords.extend(split_or_tokens(token))
                continue
            
            # 普通关键词（支持通配符 * 和 ?）
            keywords.append(token.lower())

        return keywords, filters, or_keywords, not_keywords
    
    def _parse_size(self, size_str):
        """解析大小字符串：1mb, 500kb, 10gb"""
        import re
        match = re.match(r'(\d+)(kb|mb|gb)?', size_str.lower())
        if match:
            num, unit = match.groups()
            multiplier = {'kb': 1024, 'mb': 1024**2, 'gb': 1024**3}.get(unit, 1)
            return int(num) * multiplier
        return 0
    
    def _parse_date(self, date_str):
        """解析日期字符串：2024-12-22"""
        import datetime
        dt = datetime.datetime.strptime(date_str, '%Y-%m-%d')
        return dt.timestamp()

    def _ensure_extractors(self):
        """初始化或检查外部解析器的可用性，缓存到实例属性。"""
        if getattr(self, '_extractors_initialized', False):
            return
        self._extractors_initialized = True

        self.pdf_supported = False
        self.docx_supported = False
        self.pptx_supported = False
        self.odt_supported = False
        self.pdf_extractor = None
        self.docx_extractor = None

        # PDF extractor: prefer PyPDF2, fallback to pdfminer
        try:
            import PyPDF2
            # PyPDF2 may emit PdfReadWarning for some encodings (e.g. GBK-EUC-H).
            # Suppress those warnings to avoid noisy output during parsing.
            try:
                import warnings
                from PyPDF2.errors import PdfReadWarning

                warnings.filterwarnings("ignore", category=PdfReadWarning)
            except Exception:
                # best-effort: if specific warning class not available, ignore
                try:
                    import warnings

                    warnings.filterwarnings("ignore")
                except Exception:
                    pass

            def _extract_pdf_pypdf2(path):
                try:
                    reader = PyPDF2.PdfReader(path)
                    parts = []
                    for p in reader.pages:
                        try:
                            t = p.extract_text() or ''
                        except Exception:
                            t = ''
                        if t:
                            parts.append(t)
                    return '\n'.join(parts)
                except Exception:
                    return ''

            self.pdf_supported = True
            self.pdf_extractor = _extract_pdf_pypdf2
        except Exception:
            try:
                from pdfminer.high_level import extract_text as _pdfminer_extract_text

                def _extract_pdf_pdfminer(path):
                    try:
                        return _pdfminer_extract_text(path) or ''
                    except Exception:
                        return ''

                self.pdf_supported = True
                self.pdf_extractor = _extract_pdf_pdfminer
            except Exception:
                self.pdf_supported = False

        # docx
        try:
            import docx

            def _extract_docx(path):
                try:
                    doc = docx.Document(path)
                    return '\n'.join([p.text for p in doc.paragraphs])
                except Exception:
                    return ''

            self.docx_supported = True
            self.docx_extractor = _extract_docx
        except Exception:
            self.docx_supported = False

        # pptx (python-pptx)
        try:
            import pptx

            def _extract_pptx(path):
                try:
                    prs = pptx.Presentation(path)
                    parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                parts.append(shape.text)
                    return '\n'.join(parts)
                except Exception:
                    return ''

            self.pptx_supported = True
            self.pptx_extractor = _extract_pptx
        except Exception:
            self.pptx_supported = False

        # odt (odfpy)
        try:
            from odf import text as odf_text
            from odf import opendocument

            def _extract_odt(path):
                try:
                    doc = opendocument.load(path)
                    texts = []
                    for elem in doc.getElementsByType(odf_text.P):
                        texts.append(''.join(t.data for t in elem.childNodes if getattr(t, 'data', None)))
                    return '\n'.join(texts)
                except Exception:
                    return ''

            self.odt_supported = True
            self.odt_extractor = _extract_odt
        except Exception:
            self.odt_supported = False

    def check_parsers(self):
        """返回当前解析器可用性的字典和建议的 pip 安装命令字符串。"""
        self._ensure_extractors()
        availability = {
            'pdf': bool(self.pdf_supported),
            'docx': bool(self.docx_supported),
            'pptx': bool(getattr(self, 'pptx_supported', False)),
            'odt': bool(getattr(self, 'odt_supported', False)),
        }
        # Suggest pip command (non-destructive): include common packages
        suggested = []
        if not availability['pdf']:
            suggested.append('PyPDF2')
            suggested.append('pdfminer.six')
        if not availability['docx']:
            suggested.append('python-docx')
        if not availability['pptx']:
            suggested.append('python-pptx')
        if not availability['odt']:
            suggested.append('odfpy')

        pip_cmd = 'pip install ' + ' '.join(sorted(set(suggested))) if suggested else ''
        return availability, pip_cmd

    def update_content_for_path(self, full_path, limit_size=10 * 1024 * 1024):
        """为单个文件解析并更新 `content_fts` 索引（递增更新）。"""
        if not self.conn or not getattr(self, 'has_content_fts', False):
            logger.debug('无法更新内容索引：数据库或 content_fts 不可用')
            return False

        self._ensure_extractors()

        try:
            with self.lock:
                cursor = self.conn.cursor()
                row = list(cursor.execute('SELECT id, extension, size FROM files WHERE full_path = ?', (full_path,)))
                if not row:
                    logger.debug(f'文件未在索引中: {full_path}')
                    return False
                fid, ext, sz = row[0]

            if sz and sz > limit_size:
                logger.debug(f'跳过大文件 content index 更新: {full_path}')
                return False

            ext_l = (ext or '').lower()
            text = ''
            # choose extractor based on extension
            if ext_l in {'.txt', '.md', '.py', '.csv', '.log', '.json', '.xml', '.html', '.htm', '.ini', '.cfg'}:
                try:
                    with open(full_path, 'rb') as f:
                        data = f.read()
                    try:
                        text = data.decode('utf-8')
                    except Exception:
                        try:
                            text = data.decode('gbk', errors='ignore')
                        except Exception:
                            text = data.decode('utf-8', errors='ignore')
                except Exception:
                    return False
            elif ext_l == '.pdf' and self.pdf_supported and self.pdf_extractor:
                try:
                    text = self.pdf_extractor(full_path) or ''
                except Exception:
                    return False
            elif ext_l == '.docx' and self.docx_supported and self.docx_extractor:
                try:
                    text = self.docx_extractor(full_path) or ''
                except Exception:
                    return False
            elif ext_l == '.pptx' and getattr(self, 'pptx_supported', False):
                try:
                    text = self.pptx_extractor(full_path) or ''
                except Exception:
                    return False
            elif ext_l == '.odt' and getattr(self, 'odt_supported', False):
                try:
                    text = self.odt_extractor(full_path) or ''
                except Exception:
                    return False
            else:
                # unsupported
                return False

            if not text:
                # nothing to index
                return False

            if len(text) > 1000000:
                text = text[:1000000]

            with self.lock:
                cursor = self.conn.cursor()
                try:
                    cursor.execute(
                        'INSERT OR REPLACE INTO content_fts(rowid, content, path, fileid) VALUES (?, ?, ?, ?)',
                        (fid, text, full_path, fid),
                    )
                except Exception:
                    try:
                        cursor.execute('DELETE FROM content_fts WHERE rowid = ?', (fid,))
                        cursor.execute(
                            'INSERT INTO content_fts(rowid, content, path, fileid) VALUES (?, ?, ?, ?)',
                            (fid, text, full_path, fid),
                        )
                    except Exception:
                        logger.debug(f'无法写入 content_fts: {full_path}')
                        return False
                if not HAS_APSW:
                    try:
                        self.conn.commit()
                    except Exception:
                        pass

            logger.info(f'✅ content_fts 更新: {full_path}')
            return True
        except Exception as e:
            logger.debug(f'update_content_for_path 错误: {e}')
            return False

    def build_content_index(self, allowed_exts=None, limit_size=10 * 1024 * 1024):
        """构建/更新文件内容的 FTS 索引（仅对支持的文本扩展名）。

        - `allowed_exts`: 可选的扩展名列表（如 ['.txt', '.md']），未指定则使用内置文本扩展集合。
        - `limit_size`: 跳过超过该大小（字节）的文件，默认 10MB。
        """
        if not self.conn:
            logger.warning("无法构建内容索引：数据库不可用")
            return
        if not getattr(self, 'has_content_fts', False):
            logger.warning("内容 FTS5 不可用，跳过构建内容索引")
            return

        text_exts = {'.txt', '.md', '.py', '.csv', '.log', '.json', '.xml', '.html', '.htm', '.ini', '.cfg'}
        pdf_exts = {'.pdf'}
        docx_exts = {'.docx'}

        if allowed_exts:
            allowed = set([e if e.startswith('.') else '.' + e for e in allowed_exts])
            text_exts = set(ext for ext in text_exts if ext in allowed)
            pdf_exts = set(ext for ext in pdf_exts if ext in allowed)
            docx_exts = set(ext for ext in docx_exts if ext in allowed)

        # ensure extractors are ready
        self._ensure_extractors()

        logger.info(
            f"开始构建内容索引，文本扩展: {sorted(list(text_exts))}, pdf: {bool(pdf_exts)}, docx: {bool(docx_exts)}, 大小上限: {limit_size} bytes"
        )

        # 收集候选文件（在锁内读取表）
        candidates = []
        with self.lock:
            cursor = self.conn.cursor()
            try:
                for row in cursor.execute("SELECT id, full_path, extension, size FROM files"):
                    fid, fp, ext, sz = row
                    if not ext:
                        continue
                    ext_l = ext.lower()
                    if sz and sz > limit_size:
                        continue
                    # filter by allowed extensions
                    if ext_l in text_exts or ext_l in pdf_exts or ext_l in docx_exts:
                        candidates.append((fid, fp, ext_l, sz))
            except Exception:
                logger.debug('读取 files 列表失败')

        total = len(candidates)
        logger.info(f"待解析文件数: {total}")
        try:
            # emit initial content progress (parsed, written, total, msg)
            self.content_progress_signal.emit(0, 0, total, f"收集到 {total} 个文件")
        except Exception:
            pass

        # worker to parse content for a single file
        def _parse_worker(item):
            fid, fp, ext_l, sz = item
            text = ''
            try:
                if ext_l in text_exts:
                    with open(fp, 'rb') as f:
                        data = f.read()
                    try:
                        text = data.decode('utf-8')
                    except Exception:
                        try:
                            text = data.decode('gbk', errors='ignore')
                        except Exception:
                            text = data.decode('utf-8', errors='ignore')
                elif ext_l in pdf_exts:
                    if not getattr(self, 'pdf_supported', False) or not getattr(self, 'pdf_extractor', None):
                        return None
                    text = self.pdf_extractor(fp) or ''
                elif ext_l in docx_exts:
                    if not getattr(self, 'docx_supported', False) or not getattr(self, 'docx_extractor', None):
                        return None
                    text = self.docx_extractor(fp) or ''
                else:
                    return None
            except FileNotFoundError:
                return None
            except PermissionError:
                return None
            except Exception:
                logger.debug(f'解析失败: {fp}')
                return None

            if not text:
                return None
            if len(text) > 1000000:
                text = text[:1000000]
            return (fid, fp, text)

        import concurrent.futures, multiprocessing
        max_workers = min(8, (multiprocessing.cpu_count() or 2) * 2)
        batch = []
        batch_size = 200

        parsed = 0
        written = 0
        canceled = False
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as ex:
            futures = {ex.submit(_parse_worker, item): item for item in candidates}
            for fut in concurrent.futures.as_completed(futures):
                # check for external stop request
                if getattr(self, '_stop_content_build', False):
                    logger.info('内容索引构建已被取消')
                    canceled = True
                    break
                res = None
                try:
                    res = fut.result()
                except Exception:
                    res = None
                if not res:
                    continue
                batch.append(res)
                if len(batch) >= batch_size:
                    # write batch to DB
                    with self.lock:
                        cursor = self.conn.cursor()
                        for fid, fp, text in batch:
                            try:
                                cursor.execute(
                                    "INSERT OR REPLACE INTO content_fts(rowid, content, path, fileid) VALUES (?, ?, ?, ?)",
                                    (fid, text, fp, fid),
                                )
                                written += 1
                            except Exception:
                                try:
                                    cursor.execute("DELETE FROM content_fts WHERE rowid = ?", (fid,))
                                    cursor.execute(
                                        "INSERT INTO content_fts(rowid, content, path, fileid) VALUES (?, ?, ?, ?)",
                                        (fid, text, fp, fid),
                                    )
                                except Exception:
                                    logger.debug(f"写入 content_fts 失败: {fp}")
                        try:
                            # emit content progress after write batch
                            parsed += len(batch)
                            self.content_progress_signal.emit(parsed, written, total, f"写入 {written} 条，最近: {batch[-1][1]}")
                        except Exception:
                            pass
                        if not HAS_APSW:
                            try:
                                self.conn.commit()
                            except Exception:
                                pass
                    batch = []

        # flush remaining
        if batch:
            with self.lock:
                cursor = self.conn.cursor()
                for fid, fp, text in batch:
                    try:
                        cursor.execute(
                            "INSERT OR REPLACE INTO content_fts(rowid, content, path, fileid) VALUES (?, ?, ?, ?)",
                            (fid, text, fp, fid),
                        )
                        written += 1
                    except Exception:
                        try:
                            cursor.execute("DELETE FROM content_fts WHERE rowid = ?", (fid,))
                            cursor.execute(
                                "INSERT INTO content_fts(rowid, content, path, fileid) VALUES (?, ?, ?, ?)",
                                (fid, text, fp, fid),
                            )
                        except Exception:
                            logger.debug(f"写入 content_fts 失败: {fp}")
                if not HAS_APSW:
                    try:
                        self.conn.commit()
                    except Exception:
                        pass
            try:
                parsed += len(batch)
                self.content_progress_signal.emit(parsed, written, total, f"完成写入剩余 {len(batch)} 条，最近: {batch[-1][1]}")
            except Exception:
                pass
        logger.info("内容索引构建完成")
        try:
            self.content_progress_signal.emit(parsed, written, total, "内容索引构建完成")
        except Exception:
            pass
        try:
            # emit finished/canceled signal
            self.content_build_finished_signal.emit(bool(canceled))
        except Exception:
            pass

    def stop_build_content(self):
        """请求中断正在进行的内容索引构建（线程安全标志）。"""
        self._stop_content_build = True

    def clear_stop_build(self):
        self._stop_content_build = False

    def clear_content_fts(self):
        """删除 `content_fts` 表中所有条目（用于回滚/清理）。"""
        if not self.conn or not getattr(self, 'has_content_fts', False):
            logger.warning('无法清理 content_fts：数据库或 content_fts 不可用')
            return False
        try:
            with self.lock:
                cursor = self.conn.cursor()
                cursor.execute('DELETE FROM content_fts')
                if not HAS_APSW:
                    try:
                        self.conn.commit()
                    except Exception:
                        pass
            logger.info('✅ content_fts 已清理')
            return True
        except Exception as e:
            logger.error(f'清理 content_fts 失败: {e}')
            return False


__all__ = ["IndexManager"]
