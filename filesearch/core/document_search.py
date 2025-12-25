"""
Office 文档内容搜索引擎
支持 Word (.docx), Excel (.xlsx), PowerPoint (.pptx), PDF 文件的内容搜索
"""

import os
import re
import logging
from pathlib import Path
from typing import List, Dict, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed

logger = logging.getLogger(__name__)

# 可选依赖检测
HAS_DOCX = False
HAS_OPENPYXL = False
HAS_PPTX = False
HAS_PYPDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    logger.debug("python-docx not available")

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    logger.debug("openpyxl not available")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    logger.debug("python-pptx not available")

try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    try:
        from PyPDF2 import PdfReader
        HAS_PYPDF = True
    except ImportError:
        logger.debug("pypdf/PyPDF2 not available")


class DocumentSearchEngine:
    """Office 文档搜索引擎"""
    
    # 最大文件大小 (50MB)
    MAX_FILE_SIZE = 50 * 1024 * 1024
    
    def __init__(self, max_workers: int = 4):
        self.max_workers = max_workers
    
    def extract_text_from_docx(self, file_path: str) -> Optional[str]:
        """从 Word 文档提取文本"""
        if not HAS_DOCX:
            return None
        
        try:
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            paragraphs.append(cell.text)
            
            return '\n'.join(paragraphs)
        except Exception as e:
            logger.debug(f"Error reading Word document {file_path}: {e}")
            return None
    
    def extract_text_from_xlsx(self, file_path: str) -> Optional[str]:
        """从 Excel 文档提取文本"""
        if not HAS_OPENPYXL:
            return None
        
        try:
            wb = load_workbook(file_path, read_only=True, data_only=True)
            content = []
            
            for sheet in wb.worksheets:
                content.append(f"[工作表: {sheet.title}]")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        content.append(row_text)
            
            wb.close()
            return '\n'.join(content)
        except Exception as e:
            logger.debug(f"Error reading Excel document {file_path}: {e}")
            return None
    
    def extract_text_from_pptx(self, file_path: str) -> Optional[str]:
        """从 PowerPoint 文档提取文本"""
        if not HAS_PPTX:
            return None
        
        try:
            prs = Presentation(file_path)
            content = []
            
            for i, slide in enumerate(prs.slides, 1):
                content.append(f"[幻灯片 {i}]")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text)
            
            return '\n'.join(content)
        except Exception as e:
            logger.debug(f"Error reading PowerPoint document {file_path}: {e}")
            return None
    
    def extract_text_from_pdf(self, file_path: str) -> Optional[str]:
        """从 PDF 文档提取文本"""
        if not HAS_PYPDF:
            return None
        
        try:
            reader = PdfReader(file_path)
            content = []
            
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    content.append(f"[第 {i} 页]")
                    content.append(text)
            
            return '\n'.join(content)
        except Exception as e:
            logger.debug(f"Error reading PDF document {file_path}: {e}")
            return None
    
    def extract_text(self, file_path: str) -> Optional[str]:
        """根据文件类型提取文本"""
        ext = Path(file_path).suffix.lower()
        
        # 检查文件大小
        try:
            file_size = os.path.getsize(file_path)
            if file_size > self.MAX_FILE_SIZE:
                logger.debug(f"File too large: {file_path} ({file_size} bytes)")
                return None
        except Exception as e:
            logger.debug(f"Error checking file size {file_path}: {e}")
            return None
        
        if ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif ext == '.xlsx':
            return self.extract_text_from_xlsx(file_path)
        elif ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        else:
            return None
    
    def search_in_document(
        self,
        file_path: str,
        pattern: str,
        is_regex: bool = False,
        case_sensitive: bool = False,
        context_lines: int = 2
    ) -> Optional[Dict]:
        """在单个文档中搜索"""
        content = self.extract_text(file_path)
        if not content:
            return None
        
        try:
            # 构建正则表达式
            if is_regex:
                flags = 0 if case_sensitive else re.IGNORECASE
                regex = re.compile(pattern, flags)
            else:
                escaped_pattern = re.escape(pattern)
                flags = 0 if case_sensitive else re.IGNORECASE
                regex = re.compile(escaped_pattern, flags)
            
            # 按行搜索
            lines = content.split('\n')
            matches = []
            
            for line_num, line in enumerate(lines, 1):
                if regex.search(line):
                    # 获取上下文
                    start = max(0, line_num - context_lines - 1)
                    end = min(len(lines), line_num + context_lines)
                    context = lines[start:end]
                    
                    matches.append({
                        'line_number': line_num,
                        'line_content': line.rstrip(),
                        'context': context,
                        'context_start': start + 1,
                    })
            
            if matches:
                return {
                    'file_path': file_path,
                    'match_count': len(matches),
                    'matches': matches[:50],  # 限制最多返回50个匹配
                    'file_type': Path(file_path).suffix.upper()[1:],  # DOCX, PDF, etc.
                }
            
        except re.error as e:
            logger.error(f"Invalid regex pattern: {pattern} - {e}")
        except Exception as e:
            logger.debug(f"Error searching in {file_path}: {e}")
        
        return None
    
    def search_in_documents(
        self,
        file_paths: List[str],
        pattern: str,
        is_regex: bool = False,
        case_sensitive: bool = False,
        context_lines: int = 2,
        progress_callback=None
    ) -> List[Dict]:
        """在多个文档中搜索"""
        results = []
        total = len(file_paths)
        
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            futures = {
                executor.submit(
                    self.search_in_document,
                    file_path,
                    pattern,
                    is_regex,
                    case_sensitive,
                    context_lines
                ): file_path
                for file_path in file_paths
            }
            
            for i, future in enumerate(as_completed(futures), 1):
                try:
                    result = future.result()
                    if result:
                        results.append(result)
                    
                    if progress_callback:
                        progress_callback(i, total)
                        
                except Exception as e:
                    file_path = futures[future]
                    logger.debug(f"Error processing {file_path}: {e}")
        
        # 按匹配数量排序
        results.sort(key=lambda x: x['match_count'], reverse=True)
        return results
    
    def get_supported_formats(self) -> Dict[str, bool]:
        """获取支持的格式"""
        return {
            'DOCX (Word)': HAS_DOCX,
            'XLSX (Excel)': HAS_OPENPYXL,
            'PPTX (PowerPoint)': HAS_PPTX,
            'PDF': HAS_PYPDF,
        }
    
    def is_supported(self, file_path: str) -> bool:
        """检查文件是否支持"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.docx' and HAS_DOCX:
            return True
        elif ext == '.xlsx' and HAS_OPENPYXL:
            return True
        elif ext == '.pptx' and HAS_PPTX:
            return True
        elif ext == '.pdf' and HAS_PYPDF:
            return True
        
        return False


__all__ = ['DocumentSearchEngine', 'HAS_DOCX', 'HAS_OPENPYXL', 'HAS_PPTX', 'HAS_PYPDF']
